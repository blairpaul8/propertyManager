from pptx import Presentation
import os

pptx_file = os.path.expanduser("test_slide.pptx")

output_folder = os.path.expanduser("~/Desktop/Extracted_Images")
os.makedirs(output_folder, exist_ok=True)

presentation = Presentation(pptx_file)


def get_number_from_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            if len(table.rows) >= 4 and len(table.columns) >= 2:
                return table.cell(2, 1).text.strip()
    return None


for slide_index, slide in enumerate(presentation.slides):
    file_name_number = get_number_from_table(slide)
    if not file_name_number:
        print(
            f"No number found in the 4th row, 2nd column on slide {slide_index + 1}. Skipping..."
        )
        continue

    for shape in slide.shapes:
        if shape.shape_type == 13:
            image = shape.image
            extension = image.ext if image.ext else "jpg"
            image_path = os.path.join(output_folder, f"{file_name_number}.{extension}")
            with open(image_path, "wb") as f:
                f.write(image.blob)
            print(f"Saved image to: {image_path}")

print(f"Images extracted and saved to: {output_folder}")
